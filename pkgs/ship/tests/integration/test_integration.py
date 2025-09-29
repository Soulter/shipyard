import pytest
import docker
import requests
import time
import uuid
from pathlib import Path


@pytest.fixture(scope="session")
def docker_client():
    """创建Docker客户端"""
    return docker.from_env()


@pytest.fixture(scope="session")
def ship_image(docker_client):
    """构建ship Docker镜像"""
    # 获取项目根目录
    project_root = Path(__file__).parent.parent.parent

    print(f"Building ship image from {project_root}...")

    # 构建镜像
    image, build_logs = docker_client.images.build(
        path=str(project_root), tag="ship:test", rm=True, forcerm=True
    )

    # 打印构建日志
    for log in build_logs:
        if "stream" in log:
            print(log["stream"].strip())

    yield image

    # 测试完成后清理镜像
    print("Cleaning up ship image...")
    docker_client.images.remove(image.id, force=True)


@pytest.fixture(scope="session")
def ship_container(docker_client, ship_image):
    """启动ship容器"""
    print("Starting ship container...")

    container = docker_client.containers.run(
        ship_image.id,
        ports={"8123/tcp": ("127.0.0.1", 0)},  # 随机端口映射
        detach=True,
        auto_remove=True,
    )

    # 获取映射的端口
    container.reload()
    port_mapping = container.attrs["NetworkSettings"]["Ports"]["8123/tcp"]
    if not port_mapping:
        raise RuntimeError("Failed to get port mapping")

    host_port = port_mapping[0]["HostPort"]
    base_url = f"http://127.0.0.1:{host_port}"

    # 等待服务启动
    max_retries = 30
    for i in range(max_retries):
        try:
            response = requests.get(f"{base_url}/health", timeout=5)
            if response.status_code == 200:
                print(f"Ship container is ready on port {host_port}")
                break
        except requests.exceptions.RequestException:
            pass

        time.sleep(1)
        if i == max_retries - 1:
            # 打印容器日志用于调试
            logs = container.logs().decode("utf-8")
            print(f"Container logs:\n{logs}")
            raise RuntimeError("Ship container failed to start")

    yield {"container": container, "base_url": base_url, "port": host_port}

    # 测试完成后停止容器
    print("Stopping ship container...")
    container.stop(timeout=10)


@pytest.fixture
def session_id():
    """生成唯一的session ID"""
    return str(uuid.uuid4())


@pytest.fixture
def api_headers(session_id):
    """生成API请求头"""
    return {"X-SESSION-ID": session_id}


class TestShipIntegration:
    """Ship集成测试"""

    @pytest.mark.integration
    @pytest.mark.integration
    def test_health_check(self, ship_container):
        """测试健康检查端点"""
        base_url = ship_container["base_url"]

        response = requests.get(f"{base_url}/health")
        assert response.status_code == 200

        data = response.json()
        assert data["status"] == "healthy"

    @pytest.mark.integration
    @pytest.mark.integration
    def test_root_endpoint(self, ship_container):
        """测试根端点"""
        base_url = ship_container["base_url"]

        response = requests.get(f"{base_url}/")
        assert response.status_code == 200

        data = response.json()
        assert "message" in data
        assert "Ship API is running" in data["message"]


class TestFilesystemComponent:
    """文件系统组件测试"""

    @pytest.mark.integration
    @pytest.mark.integration
    def test_create_and_read_file(self, ship_container, api_headers):
        """测试创建和读取文件"""
        base_url = ship_container["base_url"]

        # 创建文件
        create_data = {
            "path": "test_file.txt",
            "content": "Hello, World!",
            "mode": 0o644,
        }
        response = requests.post(
            f"{base_url}/fs/create_file", json=create_data, headers=api_headers
        )
        assert response.status_code == 200

        # 读取文件
        read_data = {"path": "test_file.txt"}
        response = requests.post(
            f"{base_url}/fs/read_file", json=read_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["content"] == "Hello, World!"
        assert data["path"].endswith("test_file.txt")

    @pytest.mark.integration
    @pytest.mark.integration
    def test_write_file(self, ship_container, api_headers):
        """测试写入文件"""
        base_url = ship_container["base_url"]

        # 写入文件
        write_data = {
            "path": "write_test.txt",
            "content": "Written content",
            "mode": "w",
        }
        response = requests.post(
            f"{base_url}/fs/write_file", json=write_data, headers=api_headers
        )
        assert response.status_code == 200

        # 验证文件内容
        read_data = {"path": "write_test.txt"}
        response = requests.post(
            f"{base_url}/fs/read_file", json=read_data, headers=api_headers
        )
        assert response.status_code == 200
        data = response.json()
        assert data["content"] == "Written content"

    @pytest.mark.integration
    @pytest.mark.integration
    def test_list_directory(self, ship_container, api_headers):
        """测试列出目录"""
        base_url = ship_container["base_url"]

        # 先创建一些文件
        files_to_create = ["file1.txt", "file2.py", "subdir/file3.md"]
        for file_path in files_to_create:
            create_data = {
                "path": file_path,
                "content": f"Content of {file_path}",
                "mode": 0o644,
            }
            requests.post(
                f"{base_url}/fs/create_file", json=create_data, headers=api_headers
            )

        # 列出根目录
        list_data = {"path": ".", "show_hidden": False}
        response = requests.post(
            f"{base_url}/fs/list_dir", json=list_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert "files" in data
        assert len(data["files"]) >= 3  # 至少有我们创建的文件

        # 检查文件信息
        file_names = [f["name"] for f in data["files"]]
        assert "file1.txt" in file_names
        assert "file2.py" in file_names
        assert "subdir" in file_names

    @pytest.mark.integration
    @pytest.mark.integration
    def test_delete_file(self, ship_container, api_headers):
        """测试删除文件"""
        base_url = ship_container["base_url"]

        # 先创建文件
        create_data = {
            "path": "to_delete.txt",
            "content": "This will be deleted",
            "mode": 0o644,
        }
        response = requests.post(
            f"{base_url}/fs/create_file", json=create_data, headers=api_headers
        )
        assert response.status_code == 200

        # 删除文件
        delete_data = {"path": "to_delete.txt"}
        response = requests.post(
            f"{base_url}/fs/delete_file", json=delete_data, headers=api_headers
        )
        assert response.status_code == 200

        # 验证文件已被删除
        read_data = {"path": "to_delete.txt"}
        response = requests.post(
            f"{base_url}/fs/read_file", json=read_data, headers=api_headers
        )
        assert response.status_code == 404


class TestIPythonComponent:
    """IPython组件测试"""

    @pytest.mark.integration
    @pytest.mark.integration
    def test_execute_simple_code(self, ship_container, api_headers):
        """测试执行简单Python代码"""
        base_url = ship_container["base_url"]

        # 执行简单的Python代码
        execute_data = {
            "code": 'print("Hello from IPython!")\nresult = 2 + 3\nprint(f"Result: {result}")',
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        print(response.text)
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "kernel_id" in data
        assert isinstance(data["output"], dict)

        # 检查输出内容
        output_text = data["output"].get("text", "")
        assert "Hello from IPython!" in output_text
        assert "Result: 5" in output_text

    @pytest.mark.integration
    @pytest.mark.integration
    def test_execute_with_error(self, ship_container, api_headers):
        """测试执行有错误的代码"""
        base_url = ship_container["base_url"]

        # 执行会产生错误的代码
        execute_data = {"code": "undefined_variable + 1", "timeout": 30}
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is False
        assert "error" in data
        assert "NameError" in data["error"]

    @pytest.mark.integration
    @pytest.mark.slow
    @pytest.mark.integration
    def test_matplotlib_support(self, ship_container, api_headers):
        """测试matplotlib支持和中文字体"""
        base_url = ship_container["base_url"]

        # 测试matplotlib基本功能
        execute_data = {
            "code": """
import matplotlib.pyplot as plt
import numpy as np

# 创建简单图表
x = np.linspace(0, 10, 100)
y = np.sin(x)
plt.figure(figsize=(8, 6))
plt.plot(x, y)
plt.title('正弦函数图表')  # 测试中文支持
plt.xlabel('X轴')
plt.ylabel('Y轴')
plt.grid(True)

# 保存图片而不显示
plt.savefig('test_plot.png')
plt.close()
print("Plot saved successfully!")
""",
            "timeout": 60,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True

        # 检查文件是否被创建
        response = requests.post(
            f"{base_url}/fs/list_dir", json={"path": "."}, headers=api_headers
        )
        assert response.status_code == 200
        files = response.json()["files"]
        file_names = [f["name"] for f in files]
        assert "test_plot.png" in file_names

    @pytest.mark.integration
    @pytest.mark.integration
    def test_kernel_persistence(self, ship_container, api_headers):
        """测试内核会话持久性"""
        base_url = ship_container["base_url"]

        # 第一次执行：定义变量
        execute_data = {
            "code": 'persistent_var = "Hello from persistent kernel!"',
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200
        assert response.json()["success"] is True

        # 第二次执行：使用之前定义的变量
        execute_data = {"code": "print(persistent_var)", "timeout": 30}
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "Hello from persistent kernel!" in output_text


class TestShellComponent:
    """Shell组件测试"""

    @pytest.mark.integration
    def test_execute_simple_command(self, ship_container, api_headers):
        """测试执行简单shell命令"""
        base_url = ship_container["base_url"]

        # 执行简单命令
        execute_data = {"command": 'echo "Hello from shell!"', "timeout": 10}
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert data["return_code"] == 0
        assert "Hello from shell!" in data["stdout"]

    @pytest.mark.integration
    def test_file_operations(self, ship_container, api_headers):
        """测试shell文件操作"""
        base_url = ship_container["base_url"]

        # 创建文件
        execute_data = {
            "command": 'echo "Shell created file" > shell_test.txt',
            "timeout": 10,
        }
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200
        assert response.json()["success"] is True

        # 读取文件
        execute_data = {"command": "cat shell_test.txt", "timeout": 10}
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "Shell created file" in data["stdout"]

    @pytest.mark.integration
    def test_working_directory(self, ship_container, api_headers):
        """测试工作目录"""
        base_url = ship_container["base_url"]

        # 创建子目录
        execute_data = {"command": "mkdir -p test_dir", "timeout": 10}
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200
        assert response.json()["success"] is True

        # 在子目录中执行命令
        execute_data = {
            "command": 'pwd && echo "In subdir" > test_file.txt',
            "cwd": "test_dir",
            "timeout": 10,
        }
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "test_dir" in data["stdout"]

    @pytest.mark.integration
    def test_environment_variables(self, ship_container, api_headers):
        """测试环境变量"""
        base_url = ship_container["base_url"]

        # 设置环境变量并使用
        execute_data = {
            "command": 'echo "TEST_VAR value: $TEST_VAR"',
            "env": {"TEST_VAR": "test_value_123"},
            "timeout": 10,
        }
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "test_value_123" in data["stdout"]

    @pytest.mark.integration
    def test_background_process(self, ship_container, api_headers):
        """测试后台进程"""
        base_url = ship_container["base_url"]

        # 启动后台进程
        execute_data = {
            "command": 'sleep 5 && echo "Background task completed" > bg_output.txt',
            "background": True,
            "timeout": 60,
        }
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "process_id" in data
        assert "pid" in data

        process_id = data["process_id"]

        # 检查进程状态
        response = requests.get(f"{base_url}/shell/processes", headers=api_headers)
        assert response.status_code == 200

        processes = response.json()
        process_ids = [p["process_id"] for p in processes["processes"]]
        assert process_id in process_ids

        # 等待后台任务完成
        time.sleep(6)

        # 检查输出文件
        read_data = {"path": "bg_output.txt"}
        response = requests.post(
            f"{base_url}/fs/read_file", json=read_data, headers=api_headers
        )
        assert response.status_code == 200

        file_data = response.json()
        assert "Background task completed" in file_data["content"]


class TestCrossComponentIntegration:
    """跨组件集成测试"""

    @pytest.mark.integration
    def test_python_to_shell_integration(self, ship_container, api_headers):
        """测试Python和Shell组件协作"""
        base_url = ship_container["base_url"]

        # 使用Python创建数据
        python_code = """
import json
data = {
    "name": "Integration Test",
    "values": [1, 2, 3, 4, 5],
    "timestamp": "2025-09-29"
}
with open("data.json", "w") as f:
    json.dump(data, f, indent=2)
print("Data file created")
"""

        execute_data = {"code": python_code, "timeout": 30}
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200
        assert response.json()["success"] is True

        # 使用Shell处理数据
        execute_data = {"command": 'cat data.json | jq ".name"', "timeout": 10}
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )

        # 注意：容器可能没有安装jq，所以我们使用更基本的命令
        execute_data = {"command": "cat data.json", "timeout": 10}
        response = requests.post(
            f"{base_url}/shell/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        assert "Integration Test" in data["stdout"]

    @pytest.mark.integration
    def test_filesystem_python_integration(self, ship_container, api_headers):
        """测试文件系统和Python组件协作"""
        base_url = ship_container["base_url"]

        # 使用文件系统API创建Python脚本
        script_content = """
def fibonacci(n):
    if n <= 1:
        return n
    return fibonacci(n-1) + fibonacci(n-2)

for i in range(10):
    print(f"fib({i}) = {fibonacci(i)}")
"""

        create_data = {"path": "fibonacci.py", "content": script_content, "mode": 0o644}
        response = requests.post(
            f"{base_url}/fs/create_file", json=create_data, headers=api_headers
        )
        assert response.status_code == 200

        # 使用Python执行脚本
        execute_data = {"code": 'exec(open("fibonacci.py").read())', "timeout": 30}
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "fib(0) = 0" in output_text
        assert "fib(9) = 34" in output_text


class TestDependencyLibraries:
    """依赖库测试"""

    @pytest.mark.integration
    def test_numpy_functionality(self, ship_container, api_headers):
        """测试NumPy库功能"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
import numpy as np

# 基本数组操作
arr = np.array([1, 2, 3, 4, 5])
print(f"Array: {arr}")
print(f"Sum: {np.sum(arr)}")
print(f"Mean: {np.mean(arr)}")

# 矩阵操作
matrix = np.array([[1, 2], [3, 4]])
print(f"Matrix:\\n{matrix}")
print(f"Matrix determinant: {np.linalg.det(matrix)}")

# 数学函数
x = np.linspace(0, np.pi, 5)
sin_x = np.sin(x)
print(f"Sin values: {sin_x}")
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "Array: [1 2 3 4 5]" in output_text
        assert "Sum: 15" in output_text
        assert "Mean: 3.0" in output_text

    @pytest.mark.integration
    def test_pandas_functionality(self, ship_container, api_headers):
        """测试Pandas库功能"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
import pandas as pd
import numpy as np

# 创建DataFrame
data = {
    'Name': ['Alice', 'Bob', 'Charlie', 'Diana'],
    'Age': [25, 30, 35, 28],
    'Score': [85, 92, 78, 96]
}
df = pd.DataFrame(data)
print("DataFrame:")
print(df)
print()

# 基本统计
print("Basic statistics:")
print(df.describe())
print()

# 数据筛选
high_scorers = df[df['Score'] > 85]
print("High scorers (Score > 85):")
print(high_scorers)
print()

# 保存为CSV
df.to_csv('test_data.csv', index=False)
print("Data saved to test_data.csv")

# 读取CSV
df_loaded = pd.read_csv('test_data.csv')
print("Data loaded from CSV:")
print(df_loaded.head())
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "Alice" in output_text
        assert "Bob" in output_text
        assert "Data saved to test_data.csv" in output_text

    @pytest.mark.integration
    def test_matplotlib_seaborn_integration(self, ship_container, api_headers):
        """测试matplotlib和seaborn集成"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np

# 设置seaborn样式
sns.set_style("whitegrid")

# 创建测试数据
np.random.seed(42)
data = pd.DataFrame({
    'x': np.random.normal(0, 1, 100),
    'y': np.random.normal(0, 1, 100),
    'category': np.random.choice(['A', 'B', 'C'], 100)
})

# 创建子图
fig, axes = plt.subplots(2, 2, figsize=(12, 10))

# 1. 散点图
axes[0, 0].scatter(data['x'], data['y'], alpha=0.6)
axes[0, 0].set_title('散点图')
axes[0, 0].set_xlabel('X值')
axes[0, 0].set_ylabel('Y值')

# 2. 直方图
axes[0, 1].hist(data['x'], bins=20, alpha=0.7, color='skyblue')
axes[0, 1].set_title('X值分布直方图')
axes[0, 1].set_xlabel('X值')
axes[0, 1].set_ylabel('频次')

# 3. 箱线图 (使用seaborn)
sns.boxplot(data=data, x='category', y='x', ax=axes[1, 0])
axes[1, 0].set_title('按类别的X值箱线图')

# 4. 热力图
correlation_matrix = data[['x', 'y']].corr()
sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', ax=axes[1, 1])
axes[1, 1].set_title('相关性热力图')

plt.tight_layout()
plt.savefig('comprehensive_plot.png', dpi=150, bbox_inches='tight')
plt.close()

print("综合图表已保存为 comprehensive_plot.png")
print(f"数据形状: {data.shape}")
print(f"X值均值: {data['x'].mean():.3f}")
print(f"Y值均值: {data['y'].mean():.3f}")
""",
            "timeout": 60,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "综合图表已保存为 comprehensive_plot.png" in output_text
        assert "数据形状: (100, 3)" in output_text

    @pytest.mark.integration
    def test_scikit_learn_functionality(self, ship_container, api_headers):
        """测试scikit-learn机器学习功能"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
from sklearn.datasets import make_classification
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LogisticRegression
from sklearn.metrics import accuracy_score, classification_report
from sklearn.preprocessing import StandardScaler
import numpy as np

# 生成模拟数据
X, y = make_classification(n_samples=1000, n_features=20, n_informative=10, 
                          n_redundant=10, n_clusters_per_class=1, random_state=42)

print(f"数据集形状: {X.shape}")
print(f"类别分布: {np.bincount(y)}")

# 数据分割
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

# 数据标准化
scaler = StandardScaler()
X_train_scaled = scaler.fit_transform(X_train)
X_test_scaled = scaler.transform(X_test)

# 训练逻辑回归模型
model = LogisticRegression(random_state=42)
model.fit(X_train_scaled, y_train)

# 预测和评估
y_pred = model.predict(X_test_scaled)
accuracy = accuracy_score(y_test, y_pred)

print(f"模型准确率: {accuracy:.3f}")
print("\\n分类报告:")
print(classification_report(y_test, y_pred))

# 特征重要性（回归系数）
feature_importance = np.abs(model.coef_[0])
top_features = np.argsort(feature_importance)[-5:]
print("\\n前5个重要特征的索引:")
print(top_features)
print("对应的重要性分数:")
print(feature_importance[top_features])
""",
            "timeout": 45,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "数据集形状: (1000, 20)" in output_text
        assert "模型准确率:" in output_text
        assert "分类报告:" in output_text

    @pytest.mark.integration
    def test_image_processing_libraries(self, ship_container, api_headers):
        """测试图像处理库 (OpenCV, PIL)"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
import cv2
import numpy as np
from PIL import Image, ImageDraw, ImageFont
import matplotlib.pyplot as plt

# 使用NumPy创建一个简单图像
height, width = 200, 300
img_array = np.zeros((height, width, 3), dtype=np.uint8)

# 绘制一些几何图形
cv2.rectangle(img_array, (50, 50), (150, 100), (0, 255, 0), -1)  # 绿色矩形
cv2.circle(img_array, (200, 150), 30, (255, 0, 0), -1)  # 蓝色圆形
cv2.line(img_array, (0, 0), (width, height), (0, 0, 255), 3)  # 红色对角线

# 保存OpenCV图像
cv2.imwrite('opencv_test.jpg', img_array)
print("OpenCV图像已保存为 opencv_test.jpg")

# 使用PIL处理图像
pil_img = Image.fromarray(cv2.cvtColor(img_array, cv2.COLOR_BGR2RGB))

# 创建一个新的PIL图像
pil_new = Image.new('RGB', (300, 200), 'white')
draw = ImageDraw.Draw(pil_new)

# 绘制文本和形状
draw.text((10, 10), "PIL测试图像", fill='black')
draw.rectangle([20, 50, 120, 100], outline='red', width=2)
draw.ellipse([150, 50, 250, 150], outline='blue', width=2)

pil_new.save('pil_test.png')
print("PIL图像已保存为 pil_test.png")

# 图像统计信息
print(f"OpenCV图像形状: {img_array.shape}")
print(f"PIL图像大小: {pil_img.size}")
print(f"PIL图像模式: {pil_img.mode}")

# 简单的图像处理操作
gray = cv2.cvtColor(img_array, cv2.COLOR_BGR2GRAY)
blurred = cv2.GaussianBlur(gray, (15, 15), 0)
cv2.imwrite('processed_image.jpg', blurred)
print("处理后的图像已保存为 processed_image.jpg")
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "OpenCV图像已保存为 opencv_test.jpg" in output_text
        assert "PIL图像已保存为 pil_test.png" in output_text
        assert "OpenCV图像形状: (200, 300, 3)" in output_text

    @pytest.mark.integration
    def test_document_processing_libraries(self, ship_container, api_headers):
        """测试文档处理库 (openpyxl, python-docx, python-pptx, PyPDF2)"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
# 测试 openpyxl (Excel处理)
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill

# 创建Excel文件
wb = Workbook()
ws = wb.active
ws.title = "测试工作表"

# 添加数据
data = [
    ["姓名", "年龄", "分数"],
    ["张三", 25, 85],
    ["李四", 30, 92],
    ["王五", 28, 78]
]

for row in data:
    ws.append(row)

# 设置样式
header_font = Font(bold=True, color="FFFFFF")
header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")

for cell in ws[1]:
    cell.font = header_font
    cell.fill = header_fill

wb.save("test_document.xlsx")
print("Excel文件已创建: test_document.xlsx")

# 读取Excel文件验证
from openpyxl import load_workbook
wb_read = load_workbook("test_document.xlsx")
ws_read = wb_read.active
print(f"Excel文件行数: {ws_read.max_row}")
print(f"Excel文件列数: {ws_read.max_column}")
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "Excel文件已创建: test_document.xlsx" in output_text
        assert "Excel文件行数: 4" in output_text

        # 测试python-docx (Word处理)
        execute_data = {
            "code": """
# 测试 python-docx (Word文档处理)
from docx import Document
from docx.shared import Inches

# 创建Word文档
doc = Document()
doc.add_heading('测试文档标题', 0)

# 添加段落
p = doc.add_paragraph('这是一个测试段落。')
p.add_run('这是粗体文本。').bold = True
p.add_run('这是斜体文本。').italic = True

# 添加列表
doc.add_heading('测试列表', level=1)
doc.add_paragraph('第一项', style='List Bullet')
doc.add_paragraph('第二项', style='List Bullet')
doc.add_paragraph('第三项', style='List Bullet')

# 添加表格
doc.add_heading('测试表格', level=1)
table = doc.add_table(rows=3, cols=3)
table.style = 'Table Grid'

# 填充表格数据
cells = table._cells
cells[0].text = '列1'
cells[1].text = '列2'
cells[2].text = '列3'
cells[3].text = '数据1'
cells[4].text = '数据2'
cells[5].text = '数据3'

doc.save('test_document.docx')
print("Word文档已创建: test_document.docx")

# 验证文档
doc_read = Document('test_document.docx')
paragraphs_count = len(doc_read.paragraphs)
print(f"Word文档段落数: {paragraphs_count}")
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "Word文档已创建: test_document.docx" in output_text

        # 测试python-pptx (PowerPoint处理)
        execute_data = {
            "code": """
# 测试 python-pptx (PowerPoint处理)
from pptx import Presentation
from pptx.util import Inches

# 创建PowerPoint演示文稿
prs = Presentation()

# 添加标题幻灯片
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "测试演示文稿"
subtitle.text = "使用python-pptx创建"

# 添加内容幻灯片
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = '测试内容'

tf = body_shape.text_frame
tf.text = '第一个要点'

p = tf.add_paragraph()
p.text = '第二个要点'
p.level = 1

p = tf.add_paragraph()
p.text = '第三个要点'
p.level = 2

prs.save('test_presentation.pptx')
print("PowerPoint演示文稿已创建: test_presentation.pptx")

# 验证演示文稿
prs_read = Presentation('test_presentation.pptx')
slides_count = len(prs_read.slides)
print(f"PowerPoint幻灯片数: {slides_count}")
""",
            "timeout": 30,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "PowerPoint演示文稿已创建: test_presentation.pptx" in output_text
        assert "PowerPoint幻灯片数: 2" in output_text


class TestDependencyIntegration:
    """依赖库集成测试"""

    @pytest.mark.integration
    @pytest.mark.slow
    def test_comprehensive_data_pipeline(self, ship_container, api_headers):
        """测试完整的数据处理流程"""
        base_url = ship_container["base_url"]

        execute_data = {
            "code": """
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.datasets import make_regression
from sklearn.model_selection import train_test_split
from sklearn.linear_model import LinearRegression
from sklearn.metrics import mean_squared_error, r2_score
from openpyxl import Workbook
import json

# 1. 数据生成
print("=== 步骤1: 生成模拟数据 ===")
X, y = make_regression(n_samples=200, n_features=5, noise=10, random_state=42)
feature_names = [f'特征_{i+1}' for i in range(5)]
df = pd.DataFrame(X, columns=feature_names)
df['目标值'] = y

print(f"数据集形状: {df.shape}")
print("前5行数据:")
print(df.head())

# 2. 数据探索
print("\\n=== 步骤2: 数据探索 ===")
print("描述性统计:")
print(df.describe())

# 3. 数据可视化
print("\\n=== 步骤3: 数据可视化 ===")
fig, axes = plt.subplots(2, 3, figsize=(15, 10))
fig.suptitle('数据探索可视化', fontsize=16)

# 目标变量分布
axes[0, 0].hist(df['目标值'], bins=20, alpha=0.7, color='skyblue')
axes[0, 0].set_title('目标值分布')
axes[0, 0].set_xlabel('目标值')
axes[0, 0].set_ylabel('频次')

# 特征相关性热力图
correlation_matrix = df.corr()
sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', ax=axes[0, 1])
axes[0, 1].set_title('特征相关性')

# 散点图：特征1 vs 目标值
axes[0, 2].scatter(df['特征_1'], df['目标值'], alpha=0.6)
axes[0, 2].set_title('特征1 vs 目标值')
axes[0, 2].set_xlabel('特征1')
axes[0, 2].set_ylabel('目标值')

# 箱线图
df.boxplot(column=['特征_1', '特征_2', '特征_3'], ax=axes[1, 0])
axes[1, 0].set_title('特征分布箱线图')

# 特征重要性（通过相关性）
feature_importance = correlation_matrix['目标值'].abs().sort_values(ascending=False)[1:]
axes[1, 1].bar(range(len(feature_importance)), feature_importance.values)
axes[1, 1].set_title('特征重要性')
axes[1, 1].set_xlabel('特征索引')
axes[1, 1].set_ylabel('与目标值的相关性')

# 预测vs实际散点图（稍后填充）
axes[1, 2].set_title('预测vs实际值')

plt.tight_layout()
plt.savefig('data_exploration.png', dpi=150, bbox_inches='tight')
print("数据探索图表已保存为 data_exploration.png")

# 4. 机器学习建模
print("\\n=== 步骤4: 机器学习建模 ===")
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)

model = LinearRegression()
model.fit(X_train, y_train)

y_pred = model.predict(X_test)

mse = mean_squared_error(y_test, y_pred)
r2 = r2_score(y_test, y_pred)

print(f"均方误差 (MSE): {mse:.2f}")
print(f"R² 分数: {r2:.3f}")

# 更新预测vs实际值图
axes[1, 2].clear()
axes[1, 2].scatter(y_test, y_pred, alpha=0.6)
axes[1, 2].plot([y_test.min(), y_test.max()], [y_test.min(), y_test.max()], 'r--', lw=2)
axes[1, 2].set_title('预测vs实际值')
axes[1, 2].set_xlabel('实际值')
axes[1, 2].set_ylabel('预测值')

plt.savefig('complete_analysis.png', dpi=150, bbox_inches='tight')
plt.close()

# 5. 结果导出
print("\\n=== 步骤5: 结果导出 ===")

# 导出到Excel
wb = Workbook()
ws = wb.active
ws.title = "模型结果"

# 写入标题
ws.append(["指标", "数值"])
ws.append(["均方误差", mse])
ws.append(["R² 分数", r2])
ws.append(["训练样本数", len(X_train)])
ws.append(["测试样本数", len(X_test)])

# 写入特征重要性
ws.append([])
ws.append(["特征", "系数"])
for i, coef in enumerate(model.coef_):
    ws.append([f"特征_{i+1}", coef])

wb.save("model_results.xlsx")
print("模型结果已保存到 model_results.xlsx")

# 导出为JSON
results_dict = {
    "model_metrics": {
        "mse": float(mse),
        "r2_score": float(r2),
        "train_samples": len(X_train),
        "test_samples": len(X_test)
    },
    "feature_coefficients": {
        f"特征_{i+1}": float(coef) for i, coef in enumerate(model.coef_)
    },
    "predictions_sample": {
        "actual": y_test[:5].tolist(),
        "predicted": y_pred[:5].tolist()
    }
}

with open('model_results.json', 'w', encoding='utf-8') as f:
    json.dump(results_dict, f, ensure_ascii=False, indent=2)

print("模型结果已保存到 model_results.json")

print("\\n=== 数据处理流程完成 ===")
print(f"共生成 {len(df)} 个样本")
print(f"使用 {len(feature_names)} 个特征")
print(f"最终模型R²分数: {r2:.3f}")
""",
            "timeout": 120,
        }
        response = requests.post(
            f"{base_url}/ipython/exec", json=execute_data, headers=api_headers
        )
        assert response.status_code == 200

        data = response.json()
        assert data["success"] is True
        output_text = data["output"].get("text", "")
        assert "数据集形状: (200, 6)" in output_text
        assert "均方误差 (MSE):" in output_text
        assert "R² 分数:" in output_text
        assert "数据处理流程完成" in output_text


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
